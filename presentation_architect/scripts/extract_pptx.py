import sys
import os

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

def extract_pptx_to_markdown(pptx_path, output_md_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} does not exist.")
        sys.exit(1)

    prs = Presentation(pptx_path)
    md_content = f"# Presentation Outline: {os.path.basename(pptx_path)}\n\n"

    for i, slide in enumerate(prs.slides):
        md_content += f"## Slide {i + 1}\n"
        
        # Extract title and text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                md_content += f"**Title**: {text}\n"
            else:
                md_content += f"- {text}\n"
        
        # Extract Chart Data
        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                md_content += "\n**[Chart Data Extracted]**\n"
                md_content += f"Chart Type: {chart.chart_type}\n\n"
                
                # Extract Categories
                categories = [c.label for c in chart.plots[0].categories]
                md_content += "| Category | " + " | ".join(categories) + " |\n"
                md_content += "|---| " + " | ".join(["---"] * len(categories)) + " |\n"
                
                # Extract Series
                for series in chart.series:
                    values = [str(v) for v in series.values]
                    md_content += f"| {series.name} | " + " | ".join(values) + " |\n"
                md_content += "\n"

        # Extract Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_content += f"\n**Speaker Notes**:\n> {notes}\n"

        md_content += "\n---\n\n"

    with open(output_md_path, 'w', encoding='utf-8') as f:
        f.write(md_content)
    
    print(f"Successfully extracted {len(prs.slides)} slides to {output_md_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx.py <input.pptx> <output.md>")
        sys.exit(1)
    
    extract_pptx_to_markdown(sys.argv[1], sys.argv[2])
